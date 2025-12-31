from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some basic slide layouts (0: Title, 1: Title+Content, etc.)
    # In standard themes: 0=Title, 1=Title+Content, 2=Section Header
    
    # 1. Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "台股智選系統\n(Taiwan Stock Selector)"
    subtitle.text = "打造個人的自動化選股助手\n\n報告者：[您的名字]"

    # 2. Motivation & Goal
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "專案動機與目標 (Why & What)"
    tf = content.text_frame
    tf.text = "為什麼要做這個？"
    
    p = tf.add_paragraph()
    p.text = "痛點：台股標的超過 2000 檔，逐一研究財報太耗時"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "目標：開發一個自動化工具"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "將複雜財務數據轉換為簡單的「0-10 分」評分"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "提供視覺化的網頁選股介面"
    p.level = 2

    # 3. Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "技術架構與關鍵函式庫"
    tf = content.text_frame
    tf.text = "核心語言：Python 3.10+"
    
    p = tf.add_paragraph()
    p.text = "介面 (UI)：Streamlit"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "快速建構 Web App，內建豐富互動組件 (Slider, Selectbox)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "資料處理：Pandas & NumPy"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "處理表格資料、計算財務指標 (ROE, PE, MA)"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "AI 整合：OpenAI API + Regex (備援)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "資料來源：FinMind API"
    p.level = 1

    # 4. Interface Demo Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Streamlit 介面操作"
    tf = content.text_frame
    tf.text = "如何操作這個系統？"
    
    p = tf.add_paragraph()
    p.text = "1. 啟動：streamlit run main.py"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 側邊欄導航：切換不同功能 (搜尋、篩選、排名)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. 互動式篩選：使用滑桿即時過濾股票"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. AI 智慧選股：混合式設計 (Hybrid Design)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "模式一：OpenAI 語意理解 (精準)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "模式二：Regex 關鍵字匹配 (穩健備援)"
    p.level = 2

    # 5. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "系統結構 (Project Structure)"
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "📁 main.py (程式入口)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "📁 pages/ (多頁面管理)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "1_搜尋.py, 2_篩選.py, 3_排名.py..."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "📁 src/ (核心邏輯模組)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "data_fetcher.py: 資料爬蟲與快取管理 (Offline-First)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "stock_analyzer.py: 11項指標計算與評分演算法"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "finmind_api.py: 串接外部 API"
    p.level = 1

    # 6. Scoring Logic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "核心演算法：評分系統"
    tf = content.text_frame
    tf.text = "如何將複雜數據變為 0-10 分？"
    
    p = tf.add_paragraph()
    p.text = "一般股票 (Stock) 權重："
    p.level = 1
    p = tf.add_paragraph()
    p.text = "獲利能力 (ROE) 40% + 估值 (PE) 30% + 其他 30%"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "ETF 權重："
    p.level = 1
    p = tf.add_paragraph()
    p.text = "殖利率 80% + 配息穩定度 20%"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "標準化處理："
    p.level = 1
    p = tf.add_paragraph()
    p.text = "將不同單位的指標 (%, 元, 倍) 映射到 0-100 分數區間"
    p.level = 2

    # 7. Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "遭遇困難與解決方案"
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "問題 1：API 回應慢，導致網頁卡頓"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "👉 解法：實作「離線快取機制」，預先下載並清洗資料"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "問題 2：資料缺漏 (NaN)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "👉 解法：使用 Pandas 進行資料清理 (fillna/dropna) 並增加容錯邏輯"
    p.level = 1

    # 8. Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[2]) # Section Header
    title = slide.shapes.title
    title.text = "Q & A\n感謝聆聽"

    prs.save('docs/Project_Presentation.pptx')
    print("Presentation generated successfully at docs/Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
